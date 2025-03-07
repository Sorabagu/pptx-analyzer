import os
from pptx import Presentation
from PIL import Image

def extract_text_from_pptx(file_path):
    """
    Extrait tout le texte d'un fichier .pptx.

    :param file_path: Chemin du fichier PowerPoint
    :return: Texte contenu dans le fichier
    """
    text_content = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
    except Exception as e:
        print(f"Erreur lecture {file_path} : {e}")
    
    return "\n".join(text_content)

def search_in_pptx(files_list, keyword):
    """
    Recherche un mot-clé dans une liste de fichiers .pptx.

    :param files_list: Liste de fichiers [(nom, chemin, date, taille)]
    :param keyword: Mot-clé à rechercher
    :return: Liste des fichiers contenant le mot-clé
    """
    matching_files = []

    for filename, filepath, date_modif, file_size in files_list:
        text = extract_text_from_pptx(filepath)
        if keyword.lower() in text.lower():
            matching_files.append((filename, filepath, date_modif, file_size))  # Correction ici ✅

    return matching_files

def extract_first_slide_image(file_path, output_dir="previews/"):
    """
    Extrait une image de la première diapositive d'un fichier .pptx.
    
    :param file_path: Chemin du fichier PowerPoint.
    :param output_dir: Dossier où sauvegarder l'image.
    :return: Chemin de l'image extraite ou None si pas d'image.
    """
    try:
        prs = Presentation(file_path)
        if not prs.slides:
            return None  # Aucune diapositive

        slide = prs.slides[0]  # Première diapositive
        os.makedirs(output_dir, exist_ok=True)
        img_path = os.path.join(output_dir, os.path.basename(file_path) + ".png")

        # Extraire une image intégrée
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Vérifier si c'est une image
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                return img_path

    except Exception as e:
        print(f"Erreur lors de l'extraction de l'image : {e}")

    return None  # Aucune image trouvée
